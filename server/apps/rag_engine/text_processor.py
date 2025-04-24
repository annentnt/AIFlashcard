
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import re
from typing import List, Dict, Any

class TextProcessor:
    def __init__(self, chunk_size=500, chunk_overlap=100):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    def extract_text_from_file(self, file, filename: str) -> str:
        """Extract text content from various file types"""
        text = ''
        if filename.endswith('.txt'):
            text = file.read().decode('utf-8')

        elif filename.endswith('.docx'):
            doc = Document(file)
            text = '\n'.join([para.text for para in doc.paragraphs])

        elif filename.endswith('.pptx'):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + '\n'

        elif filename.endswith('.pdf'):
            pdf_file = fitz.open(stream=file.read(), filetype="pdf")
            for page in pdf_file:
                text += page.get_text()

        return text
    
    def chunk_text(self, text: str) -> List[str]:
        """Split text into overlapping chunks for better processing"""
        chunks = []
        
        # Simple paragraph splitting
        paragraphs = re.split(r'\n\s*\n', text)
        current_chunk = ""
        
        for paragraph in paragraphs:
            # Clean the paragraph
            clean_paragraph = paragraph.strip()
            if not clean_paragraph:
                continue
                
            # If adding this paragraph exceeds chunk size, save current chunk and start new one
            if len(current_chunk) + len(clean_paragraph) > self.chunk_size:
                if current_chunk:
                    chunks.append(current_chunk)
                
                # For very long paragraphs that exceed chunk size on their own
                if len(clean_paragraph) > self.chunk_size:
                    # Split by sentences or just hard-split if needed
                    sentences = re.split(r'(?<=[.!?])\s+', clean_paragraph)
                    current_chunk = ""
                    
                    for sentence in sentences:
                        if len(current_chunk) + len(sentence) > self.chunk_size:
                            if current_chunk:
                                chunks.append(current_chunk)
                            current_chunk = sentence + " "
                        else:
                            current_chunk += sentence + " "
                else:
                    current_chunk = clean_paragraph + "\n\n"
            else:
                current_chunk += clean_paragraph + "\n\n"
        
        # Add the last chunk if it has content
        if current_chunk:
            chunks.append(current_chunk)
            
        # Create overlapping chunks for better context preservation
        overlapping_chunks = []
        for i in range(len(chunks)):
            # If not the first chunk, add some context from the previous chunk
            if i > 0:
                prev_chunk_end = chunks[i-1][-self.chunk_overlap:]
                current_chunk = prev_chunk_end + chunks[i]
            else:
                current_chunk = chunks[i]
                
            # If not the last chunk, add some context from the next chunk
            if i < len(chunks) - 1:
                next_chunk_start = chunks[i+1][:self.chunk_overlap]
                current_chunk = current_chunk + next_chunk_start
                
            overlapping_chunks.append(current_chunk)
            
        return overlapping_chunks





