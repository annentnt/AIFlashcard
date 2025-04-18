# utils.py
import openai
import json
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF

from memoria.settings import OPENAI_API_KEY

def is_content_safe(text: str) -> bool:
    moderation_client = openai.OpenAI(api_key=OPENAI_API_KEY)

    moderation_response = moderation_client.moderations.create(
        model="omni-moderation-latest",
        input=text
    )

    result = moderation_response.results[0].flagged

    return result

def extract_text_from_file(file, filename):
    text = ''
    if filename.endswith('.txt'):
        text = file.read().decode('utf-8')

    elif filename.endswith('.docx'):
        doc = Document(file)
        text = '\n'.join([para.text for para in doc.paragraphs])

    elif filename.endswith('.pptx'):
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + '\n'

    elif filename.endswith('.pdf'):
        pdf_file = fitz.open(stream=file.read(), filetype="pdf")
        for page in pdf_file:
            text += page.get_text()

    # print("Extracted text:", text)

    return text

def call_openai_extract_flashcards(file, filename, num_flashcards):
    text = extract_text_from_file(file, filename)

    if is_content_safe(text):
        return None

    prompt = f"""
        You are a language assistant that helps learners understand new vocabulary from documents. First, check if the text contains violent, hateful, or inappropriate content. If yes, just respond with a blank JSON object.

        Otherwise, extract up to {num_flashcards} key vocabulary words or phrases using Named Entity Recognition (NER) and keyphrase extraction. For each item, provide a description that fits the context of the original text.

        Return the result as a JSON schema
    """

    client = openai.OpenAI(api_key=OPENAI_API_KEY)

    response = client.responses.create(
        model="gpt-4o-mini",
        input=[
            {"role": "system", "content": prompt},
            {"role": "user", "content": text}
        ],
        temperature=0.3,
        text={
            "format": {
                "type": "json_schema",
                "name": "vocab_list",
                "schema": {
                    "type": "object",
                    "properties": {
                        "vocabs": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "vocabulary": {"type": "string"},
                                    "description": {"type": "string"}
                                },
                                "required": ["vocabulary", "description"],
                                "additionalProperties": False
                            }
                        }
                    },
                    "required": ["vocabs"],
                    "additionalProperties": False
                },
                "strict": True
            }
        }
    )

    output = json.loads(response.output_text)

    return output["vocabs"]